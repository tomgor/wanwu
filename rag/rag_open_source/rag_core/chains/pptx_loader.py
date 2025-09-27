from pptx import Presentation
from typing import List, Optional

from langchain_core.documents import Document
from langchain_community.document_loaders import TextLoader

import os
import nltk
current_file_path = os.path.abspath(__file__)
# 获取当前文件所在的目录
current_dir = os.path.dirname(current_file_path)
# 添加项目根目录到 sys.path
# sys.path.append(root_dir)
# 拼接nltk_data文件夹的路径
nltk_data_path = os.path.join(current_dir, 'nltk_data')
nltk.data.path.append(nltk_data_path)
# import minio_utils
from logging_config import setup_logging
logger_name = 'rag_pptx_loader'
app_name = os.getenv("LOG_FILE")
logger = setup_logging(app_name,logger_name)
logger.info(logger_name+'---------LOG_FILE：'+repr(app_name))


def table_convert_html(table):
    # embedding_content = []
    #
    # table_title_list = []

    def is_empty(cell):
        return cell is None
        # return cell is None or str(cell).strip() == ''

    def get_colspan(row, col_idx, processed_cells, row_idx):
        # 计算当前单元格的跨列数
        if (row_idx, col_idx) in processed_cells:
            return 0  # 跳过已处理的单元格
        # 计算当前单元格的跨列数
        colspan = 1
        for i in range(col_idx + 1, len(row)):
            if is_empty(row[i]) and (row_idx, i) not in processed_cells:
                processed_cells.add((row_idx, i))  # 标记已处理
                colspan += 1
            else:
                break
        return colspan

    def get_rowspan(rows, row_idx, col_idx, processed_cells):
        if (row_idx, col_idx) in processed_cells:
            return 0  # 跳过已处理的单元格
        rowspan = 1
        for i in range(row_idx + 1, len(rows)):
            if is_empty(rows[i][col_idx]) and (i, col_idx) not in processed_cells:
                processed_cells.add((i, col_idx))  # 标记已处理
                rowspan += 1
            else:
                break
        return rowspan

    # 开始构建 HTML 表格
    html = "<table border='1'>\n"
    processed_cells = set()  # 记录已处理的单元格位置 (row_idx, col_idx)
    for row_idx, row in enumerate(table):
        html += "<tr>"
        row_text = "<tr>"
        col_idx = 0
        while col_idx < len(row):
            if (row_idx, col_idx) in processed_cells:
                col_idx += 1  # 跳过已处理的单元格，避免死循环
                continue
            cell = row[col_idx]

            if is_empty(cell):
                col_idx += 1
                continue  # 跳过空单元格

            # 获取跨列数
            colspan = get_colspan(row, col_idx, processed_cells, row_idx)

            # 获取跨行数
            rowspan = get_rowspan(table, row_idx, col_idx, processed_cells)

            # 添加单元格
            if colspan > 1 or rowspan > 1:
                html += f"<td colspan='{colspan}' rowspan='{rowspan}'>{cell}</td>"
                row_text += f"<td colspan='{colspan}' rowspan='{rowspan}'>{cell}</td>"
            else:
                html += f"<td>{cell}</td>"
                row_text += f"<td>{cell}</td>"
            # 把除首行外的前两列的内容所谓embedding索引
            # if row_idx > 0 and col_idx < 2 and str(cell).strip() != '':
            #     # 对于表格中没有行分割线的表：通过 换行符切割索引
            #     if len(table) <= 2 and "\n" in str(cell):
            #         cell_list = str(cell).split("\n")
            #         embedding_content.extend(cell_list)
            #     else:
            #         embedding_content.append(cell)
            # 标记已经处理过的单元格
            for i in range(rowspan):
                for j in range(colspan):
                    processed_cells.add((row_idx + i, col_idx + j))
            # 跳过已经处理过的跨列单元格
            col_idx += colspan

        html += "</tr>\n"
        row_text += "</tr>"
        # if row_idx > 0:
        #     embedding_content.append(row_text)

    # 结束 HTML 表格
    html += "</table>\n"
    # print(json.dumps(embedding_content,ensure_ascii=False))
    return html


class PPTXLoader(TextLoader):
    def load(self) -> List[Document]:
        text = ""
        try:
            prs = Presentation(self.file_path)
            print(prs)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        t = text_frame.text
                        text += t + '\n'
                    elif shape.has_table:
                        one_table_data = []
                        for row in shape.table.rows:  # 读每行
                            row_data = []
                            for cell in row.cells:  # 读一行中的所有单元格
                                if cell.text != "":
                                    row_data.append(cell.text)
                                else:
                                    row_data.append(None)
                                # cell.text = cell.text if cell.text != "" else ""
                                # c = cell.text
                                # row_data.append(c)
                            one_table_data.append(row_data)  # 把每一行存入表

                        print("one_table_data=%s" % one_table_data)
                        table_html = table_convert_html(one_table_data)
                        text += table_html + '\n'
        except Exception as e:
            raise RuntimeError(f"Error loading {self.file_path}") from e

        metadata = {"source": self.file_path}
        return [Document(page_content=text, metadata=metadata)]
if __name__ == "__main__":

    filepath = "./your_file.pptx"
    loader = PPTXLoader(filepath)
    docs = loader.load()
    for doc in docs:
        print(doc)