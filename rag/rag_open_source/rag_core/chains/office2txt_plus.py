import os
from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextContainer
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor, as_completed
from tqdm import tqdm

def extract_text_from_doc(doc_file_path, output_file_path):
    # 仅支持ms office 2007及以后保存的.doc文档\
    # 安装依赖：apt-get install catdoc
    full_text = os.popen("catdoc -s='utf-8' -d='utf-8' " + doc_file_path).read()
    with open(output_file_path, "w", encoding='utf-8') as file:
        file.write(full_text)
    # print(f"PDF文件 {doc_file_path} 已成功处理。")

def extract_text_from_pdf(pdf_file_path, output_file_path):
    full_text = ''
    for i, page_layout in enumerate(extract_pages(pdf_file_path)):
        for element in page_layout:
            if isinstance(element, LTTextContainer):
                full_text += element.get_text() + '\n'
    
    with open(output_file_path, "w", encoding='utf-8') as file:
        file.write(full_text)
    # print(f"PDF文件 {pdf_file_path} 已成功处理。")

def extract_text_from_pptx(ppt_file_path, output_file_path):
    # 安装依赖：pip install python-pptx
    try:
        prs = Presentation(ppt_file_path)
    except Exception as e:
        print(f"无法打开 PowerPoint 文件：{e}")
        return

    full_text = ''
    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    text = text_frame.text
                    if text:
                        full_text += text + '\n'
    except Exception as e:
        print(f"处理 PowerPoint 文件时出错：{e}")
        return
                    
    with open(output_file_path, "w", encoding='utf-8') as file:
        file.write(full_text)
    # print(f"PPT文件 {ppt_file_path} 已成功处理。")

def process_file(file_path, output_folder, file_format, progress):
    output_file_path = os.path.join(output_folder, os.path.basename(file_path).replace(file_format, '.txt'))
    if file_format == '.pptx':
        extract_text_from_pptx(file_path, output_file_path)
    elif file_format == '.pdf':
        extract_text_from_pdf(file_path, output_file_path)
    elif file_format == '.doc':
        extract_text_from_doc(file_path, output_file_path)
    progress.update(1)

def batch_extract_text(source_folder, output_folder, file_format, max_workers=10):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
        
    files = [os.path.join(source_folder, f) for f in os.listdir(source_folder) if f.endswith(file_format)]
    count = len(files)

    with ThreadPoolExecutor(max_workers=max_workers) as executor, tqdm(total=count) as progress:
        futures = [executor.submit(process_file, file_path, output_folder, file_format, progress) for file_path in files]
        for future in as_completed(futures):
            pass


